"""tc-runner 사내 소개 PPT 생성기 (12 slides, 전사 혼합 청중).

산출: doc/tc_runner_intro.pptx
강조 축: AI 활용 방식 / 실제 적용 결과물 / 비용·시간 절감 효과

기존 gen_seniorshield_ppt.py 패턴 재활용.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(BASE, "doc")
os.makedirs(OUT_DIR, exist_ok=True)
OUT = os.path.join(OUT_DIR, "tc_runner_intro.pptx")

# ── 색상 팔레트 ─────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
BLUE = RGBColor(0x29, 0x80, 0xB9)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xFA)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
PURPLE_BG = RGBColor(0xE8, 0xDA, 0xEF)
SOFT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
SOFT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
SOFT_ORANGE = RGBColor(0xFA, 0xE5, 0xD3)
SOFT_RED = RGBColor(0xFA, 0xDB, 0xD8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


# ── 헬퍼 ────────────────────────────────────────────────
def add_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, text, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                   SW, Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.20)


def add_text_box(slide, left, top, width, height, text, size=18, bold=False,
                 color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, size=16,
                    color=NAVY, line_space=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(line_space)
    return tf


def add_card(slide, left, top, width, height, title, body_lines,
             title_color=NAVY, body_color=NAVY, bg=WHITE,
             title_size=18, body_size=14):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                  width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.color.rgb = GRAY
    rect.line.width = Pt(0.75)
    # remove default text
    rect.text_frame.text = ""
    # title
    add_text_box(slide, left + Inches(0.18), top + Inches(0.10),
                 width - Inches(0.36), Inches(0.5),
                 title, size=title_size, bold=True, color=title_color)
    # body
    add_bullet_list(slide, left + Inches(0.18),
                    top + Inches(0.55),
                    width - Inches(0.36), height - Inches(0.65),
                    body_lines, size=body_size, color=body_color)


def add_arrow(slide, left, top, width, height, color=BLUE):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top,
                                   width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    arrow.text_frame.text = ""


# ── Slide 1 — Title ─────────────────────────────────────
def slide_title():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, NAVY)
    add_text_box(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.4),
                 "tc-runner", size=72, bold=True, color=WHITE)
    add_text_box(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.8),
                 "사내 AI 활용 사례 — Android QA 자동화 + 누적 학습 루프",
                 size=28, color=WHITE)
    add_text_box(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.5),
                 "발표자: 진우  ·  altech QA  ·  2026-05",
                 size=18, color=SOFT_BLUE)
    add_text_box(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
                 "Powered by Claude Code (Opus 4.7 + Sonnet 4.6)",
                 size=14, color=GRAY)


# ── Slide 2 — 문제 정의 ─────────────────────────────────
def slide_problem():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title_bar(s, "왜 만들었나 — Android QA의 반복 비용")

    add_text_box(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
                 "기존 수동 QA 흐름의 3가지 누수",
                 size=22, bold=True, color=NAVY)

    cards = [
        ("매번 다시 푼다", [
            "신규 앱 출시마다 화면·메뉴·셀렉터를",
            "처음부터 다시 탐색",
            "선임의 머릿속 지식이 휘발",
            "후임은 항상 zero-base"
        ], SOFT_RED, RED),
        ("회귀가 비싸다", [
            "수정 빌드 1회 = 수십개 시나리오 재실행",
            "수동 클릭 반복 → 누적 피로",
            "발견 못한 회귀가 출시 후에야 표면화"
        ], SOFT_ORANGE, ORANGE),
        ("데이터가 안 남는다", [
            "테스트 결과는 엑셀/메신저에만 남음",
            "다음 주기에 재사용 불가",
            "BUG·SPEC_GAP 트래킹 산발",
            "“왜 이렇게 했지?”가 영구 미해결"
        ], SOFT_BLUE, BLUE),
    ]
    x = Inches(0.7)
    for title, body, bg, tcolor in cards:
        add_card(s, x, Inches(2.2), Inches(4.0), Inches(4.0),
                 title, body, title_color=tcolor, bg=bg,
                 title_size=20, body_size=14)
        x += Inches(4.2)

    add_text_box(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
                 "→ tc-runner는 이 3가지를 동시에 해결하는 “누적 학습 루프”다.",
                 size=18, bold=True, color=NAVY)


# ── Slide 3 — 한 줄 + 3축 가치 ──────────────────────────
def slide_overview():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title_bar(s, "tc-runner는 단순 자동화가 아니다")

    add_text_box(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.7),
                 "“테스트 케이스를 사람이 짜고, AI가 실행하고, 데이터가 누적된다.”",
                 size=24, bold=True, color=NAVY)
    add_text_box(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
                 "단순 실행기가 아닌 — 신규 앱 빠른 파악 / 수정판 delta 재탐색 / 누적 데이터 보존을 동시에 수행하는 학습 루프.",
                 size=14, color=GRAY)

    cards = [
        ("AI 활용 방식", [
            "Claude Code 기반 multi-agent 워크플로우",
            "Opus가 계획, Sonnet이 실행",
            "프롬프트 자산 (STAGE1/STAGE2)",
            "agent 위임으로 컨텍스트 분리",
        ], SOFT_BLUE, BLUE),
        ("실제 적용 결과물", [
            "Music · Gallery · MiniFile",
            "Settings · SeniorShield · BUG 트래킹",
            "단말×앱 7+ 조합, TC 50+",
            "BUG·SPEC_GAP 누적 발견",
        ], SOFT_GREEN, GREEN),
        ("비용·시간 절감", [
            "TC 1건 수동 30분 → AI 컴파일 5분",
            "회귀 자동 runtime, 단계당 수십초",
            "신규 앱 진입 시간 단축",
            "데이터 누적으로 다음 주기 가속",
        ], SOFT_ORANGE, ORANGE),
    ]
    x = Inches(0.7)
    for title, body, bg, tcolor in cards:
        add_card(s, x, Inches(3.2), Inches(4.0), Inches(3.6),
                 title, body, title_color=tcolor, bg=bg,
                 title_size=22, body_size=15)
        x += Inches(4.2)


# ── Slide 4 — AI 활용 (1) Multi-agent ───────────────────
def slide_ai_agents():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title_bar(s, "AI 활용 ① — 역할 분리한 Multi-agent")

    add_text_box(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.5),
                 "사람과 AI가 단순 1:1로 묻고 답하지 않는다. 모델별 역할을 분리한다.",
                 size=18, color=NAVY)

    # Two big cards
    add_card(s, Inches(0.7), Inches(2.0), Inches(6.0), Inches(4.5),
             "🧠 Opus 4.7 — 계획 (Architect)", [
                 "복잡한 컨텍스트 종합",
                 "PR scope·plan 문서 작성",
                 "테스트 전략·리스크 분석",
                 "GATE 게이트 체인 설계",
                 "대화 길어져도 일관성 유지",
                 "예: PR 6/7 scope·plan 문서 5+ 작성",
             ], title_color=BLUE, bg=SOFT_BLUE,
             title_size=22, body_size=16)
    add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.5),
             "⚡ Sonnet 4.6 — 실행 (Builder)", [
                 "구체적 코드·스크립트 생성",
                 "TC YAML 작성, 좌표 채움",
                 "validate / runtime 실행",
                 "리포트 자동 생성",
                 "응답 속도 빠름 → 반복 작업 적합",
                 "예: SMOKE TC 30+ 구현 실행",
             ], title_color=ORANGE, bg=SOFT_ORANGE,
             title_size=22, body_size=16)

    add_text_box(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
                 "→ 같은 프로젝트도 “계획자 / 실행자”를 나누면, 컨텍스트 휘발 없이 깊은 작업이 가능하다.",
                 size=16, bold=True, color=NAVY)


# ── Slide 5 — AI 활용 (2) Prompt 자산 + 누적 루프 ────────
def slide_ai_loop():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title_bar(s, "AI 활용 ② — Prompt 자산 + 누적 학습 루프")

    add_text_box(s, Inches(0.7), Inches(1.2), Inches(6.0), Inches(0.5),
                 "Prompt = 코드처럼 버전 관리하는 자산",
                 size=20, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.7), Inches(1.8), Inches(6.0), Inches(2.5),
                    [
                        "STAGE1_NORMALIZE.md — 원본 → CTF 정규화 규칙",
                        "STAGE2_COMPILE.md — CTF → 실행 TC 컴파일",
                        "OPERATIONAL_RULES.md — 운영 규칙 항상 적용",
                        "device_profile.yaml / runner_capability.yaml",
                        "→ 매번 다시 설명 X. 모델 바뀌어도 동일 결과.",
                    ], size=15)

    add_text_box(s, Inches(0.7), Inches(4.4), Inches(6.0), Inches(0.5),
                 "누적 학습 루프 (catalog · delta)",
                 size=20, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.7), Inches(5.0), Inches(6.0), Inches(2.0),
                    [
                        "탐색 결과 = catalog/screens.json 누적",
                        "수정 빌드 → delta 비교 (jaccard / sha256)",
                        "재탐색 비용 ↓, 회귀 발견율 ↑",
                        "“다음 작업”이 “이전 데이터”에 올라탐",
                    ], size=15)

    # Right side — diagram
    box_w, box_h = Inches(2.4), Inches(0.7)
    base_x = Inches(8.5)
    items = [
        ("탐색 / probe", BLUE),
        ("catalog 누적", GREEN),
        ("TC 작성 (AI)", ORANGE),
        ("runtime 실행", BLUE),
        ("delta 분석", PURPLE_BG),
        ("학습 ← 다음 작업", GREEN),
    ]
    y = Inches(1.6)
    for i, (label, c) in enumerate(items):
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, base_x, y,
                                  box_w, box_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = c
        rect.line.color.rgb = NAVY
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        y += Inches(0.85)
        if i < len(items) - 1:
            arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                       base_x + Inches(1.0), y - Inches(0.18),
                                       Inches(0.4), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()


# ── Slide 6 — 파이프라인 다이어그램 ─────────────────────
def slide_pipeline():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title_bar(s, "파이프라인 한눈에 — 입력에서 누적까지")

    stages = [
        ("원본 TC\n(엑셀/메모)", BLUE),
        ("Stage 1\n정규화 (CTF)", GREEN),
        ("Stage 2\n컴파일 (YAML)", ORANGE),
        ("validate_tc.py\nlint+schema", BLUE),
        ("cli run\n실 단말 실행", GREEN),
        ("리포트+catalog\n누적", ORANGE),
    ]
    box_w, box_h = Inches(1.85), Inches(1.2)
    gap = Inches(0.18)
    total_w = box_w * len(stages) + gap * (len(stages) - 1)
    start_x = (SW - total_w) / 2
    y = Inches(2.4)

    for i, (label, c) in enumerate(stages):
        x = start_x + (box_w + gap) * i
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                                  box_w, box_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = c
        rect.line.color.rgb = NAVY
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        if i < len(stages) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + box_w + Inches(0.005),
                y + box_h / 2 - Inches(0.12),
                gap - Inches(0.01), Inches(0.24))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()

    # Below — annotations
    annotations = [
        ("AI가 닿는 구간", Inches(0.7), Inches(4.0), Inches(6.0),
         "Stage 1·2는 Claude가 자동 수행. 사람은 검토·승인만."),
        ("기계 검증 + 실 단말 검증", Inches(7.0), Inches(4.0), Inches(5.6),
         "validate(정적) → cli run(실 단말) 두 단계 게이트."),
        ("데이터가 남는다", Inches(0.7), Inches(5.4), Inches(6.0),
         "catalog/ + reports/ 에 누적. delta·회귀 분석에 재사용."),
        ("GATE 1→2→3→4", Inches(7.0), Inches(5.4), Inches(5.6),
         "각 단계 통과 강제. PASS 어휘 4종 분리(validate/runtime/manual/BUG)."),
    ]
    for title, x, y, w, body in annotations:
        add_text_box(s, x, y, w, Inches(0.4),
                     "▶ " + title, size=16, bold=True, color=NAVY)
        add_text_box(s, x + Inches(0.3), y + Inches(0.4),
                     w - Inches(0.3), Inches(0.6),
                     body, size=13, color=GRAY)


# ── Slide 7 — 사례 1 Music ──────────────────────────────
def slide_case_music():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title_bar(s, "적용 사례 ① — Music (com.mive.music · ODIN2)")

    add_text_box(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.5),
                 "신규 앱 cold-start로 6개 SMOKE 시나리오 자동화. 모두 실 단말 PASS.",
                 size=18, color=NAVY)

    # SMOKE 표
    rows = [
        ("SMOKE_01", "앱 콜드 런치 + HOME 라벨 6종", "10/10 step", "18.5s"),
        ("SMOKE_02", "HOME 4 탭 네비게이션 (좌표 tap 0)", "24/24 step", "50.6s"),
        ("SMOKE_03", "첫 곡 재생 진입", "12/12 step", "—"),
        ("SMOKE_04", "검색 입력 / 키보드 포커스", "PASS", "—"),
        ("SMOKE_05", "검색 결과 화면", "PASS", "—"),
        ("SMOKE_06", "즐겨찾기 add/remove + cleanup", "PASS", "—"),
    ]
    table_top = Inches(2.0)
    row_h = Inches(0.55)
    cols = [
        ("SMOKE", Inches(1.4)),
        ("시나리오", Inches(6.8)),
        ("Runtime", Inches(2.0)),
        ("Time", Inches(1.5)),
    ]
    # header
    x = Inches(0.7)
    for col, w in cols:
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_top, w, row_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = NAVY
        rect.line.color.rgb = NAVY
        tf = rect.text_frame
        p = tf.paragraphs[0]
        p.text = col
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        x += w
    # rows
    for ri, row in enumerate(rows):
        x = Inches(0.7)
        bg = WHITE if ri % 2 == 0 else SOFT_BLUE
        y = table_top + row_h * (ri + 1)
        for (col, w), val in zip(cols, row):
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, row_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = bg
            rect.line.color.rgb = GRAY
            tf = rect.text_frame
            p = tf.paragraphs[0]
            p.text = val
            p.font.size = Pt(13)
            p.font.color.rgb = NAVY
            p.alignment = PP_ALIGN.CENTER
            x += w

    add_text_box(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.4),
                 "✓ 핵심: 좌표 tap 0건 (모두 텍스트/anchor 기반) — 화면 변경에도 견고",
                 size=16, bold=True, color=GREEN)
    add_text_box(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
                 "✓ catalog 자동 누적 — 화면 baseline·delta 검출 기반 마련",
                 size=16, bold=True, color=GREEN)


# ── Slide 8 — 사례 2 종합 ───────────────────────────────
def slide_case_others():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title_bar(s, "적용 사례 ② — 단말×앱 7개 트랙 누적")

    cards = [
        ("Gallery (자체 앱)", [
            "23 PASS + 6 BUG/GAP",
            "권한·휴지통·즐겨찾기 시나리오",
            "preset 자동 주입 스크립트",
        ], SOFT_GREEN, GREEN),
        ("MiniFile (자체 앱)", [
            "26/28 PASS",
            "Opus 계획 + Sonnet 실행 분리 사례",
            "viewer intent 2건만 개발팀 대기",
        ], SOFT_BLUE, BLUE),
        ("Settings (시스템 앱)", [
            "SMOKE_01 11/11 step PASS (오늘)",
            "ROOT 6 anchor verify + cleanup",
            "신규 앱 진입 → PASS 1시간 이내",
        ], SOFT_ORANGE, ORANGE),
        ("SeniorShield (사내 앱)", [
            "11/11 PASS + A/B/C 회귀 4/4",
            "팝업 snooze 이슈 직접 수정",
            "쿨다운/세션 상태머신 분석",
        ], PURPLE_BG, NAVY),
        ("BUG 트래킹", [
            "BUG-5426 APN 모니터링",
            "BUG-17126 WCDMA Reject",
            "BUG-21838 MCFG·apns-conf 경로",
        ], SOFT_RED, RED),
        ("도구 누적", [
            "validate_tc.py + lint",
            "git_safe_push_audit.py (PR 6)",
            "synthetic_delta_measure.py (PR 7A)",
        ], SOFT_BLUE, BLUE),
    ]
    cols = 3
    cw, ch = Inches(4.0), Inches(2.55)
    gap_x, gap_y = Inches(0.15), Inches(0.20)
    start_x = Inches(0.7)
    start_y = Inches(1.4)
    for i, (title, body, bg, tc) in enumerate(cards):
        r, c = divmod(i, cols)
        x = start_x + (cw + gap_x) * c
        y = start_y + (ch + gap_y) * r
        add_card(s, x, y, cw, ch, title, body,
                 title_color=tc, bg=bg, title_size=18, body_size=13)

    add_text_box(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
                 "→ 데이터·도구·BUG·문서가 모두 한 repo에 누적된다.",
                 size=16, bold=True, color=NAVY)


# ── Slide 9 — 비용·시간 (1) ─────────────────────────────
def slide_cost_one():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title_bar(s, "비용·시간 절감 ① — TC 1건 작성 비교")

    add_text_box(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.5),
                 "수동 작성 → AI 컴파일로 단계별 시간 단축",
                 size=20, bold=True, color=NAVY)

    rows = [
        ("단계", "수동", "AI 컴파일", "절감"),
        ("원본 → 정규화 (Stage 1)", "20~30분", "30초~1분", "≈ 95% ↓"),
        ("정규화 → YAML 컴파일 (Stage 2)", "30~40분", "1~2분", "≈ 95% ↓"),
        ("validate / lint 통과", "수동 검토 10분", "1초 + 자동 lint", "≈ 99% ↓"),
        ("실 단말 runtime", "수동 클릭 5~10분", "20~50초 자동", "≈ 80% ↓"),
        ("리포트 작성", "엑셀 수동 10분", "gen_excel 자동", "≈ 95% ↓"),
        ("총 1건당", "70~100분", "5~10분", "≈ 90% ↓"),
    ]
    table_top = Inches(2.0)
    row_h = Inches(0.55)
    cols = [Inches(4.6), Inches(2.6), Inches(2.6), Inches(2.2)]
    total_w = sum(cols, Inches(0))
    start_x = (SW - total_w) / 2

    for ri, row in enumerate(rows):
        x = start_x
        is_header = ri == 0
        is_total = ri == len(rows) - 1
        bg_row = NAVY if is_header else (SOFT_GREEN if is_total else
                                         (WHITE if ri % 2 == 1 else SOFT_BLUE))
        text_color = WHITE if is_header else NAVY
        y = table_top + row_h * ri
        for w, val in zip(cols, row):
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, row_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = bg_row
            rect.line.color.rgb = GRAY
            tf = rect.text_frame
            p = tf.paragraphs[0]
            p.text = val
            p.font.size = Pt(14 if not is_total else 15)
            p.font.bold = is_header or is_total
            p.font.color.rgb = text_color
            p.alignment = PP_ALIGN.CENTER if w != cols[0] else PP_ALIGN.LEFT
            tf.margin_left = Inches(0.15)
            x += w

    add_text_box(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
                 "※ 추정치 — 단일 TC 기준. 복잡 시나리오일수록 절감 폭 ↑",
                 size=13, color=GRAY)


# ── Slide 10 — 비용·시간 (2) 누적 가속 ──────────────────
def slide_cost_two():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title_bar(s, "비용·시간 절감 ② — 누적 학습으로 가속")

    add_text_box(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.5),
                 "“N+1번째 앱은 N번째보다 빠르다” — 데이터·도구·prompt가 누적되기 때문.",
                 size=18, bold=True, color=NAVY)

    rows = [
        ("앱 도입 단계", "1번째 앱 (Music)", "최근 앱 (Settings)"),
        ("Phase 0 (탐색)", "1~2일 — 메뉴트리 직접 탐색", "30분 — probe + texts 추출 자동"),
        ("SMOKE_01 작성", "수일 — anchor·구조 수동 결정", "30분 — probe → scope → YAML"),
        ("validate / runtime 도달", "다음 날 — 디버깅 반복", "당일 — validate→runtime 1회 PASS"),
        ("회귀 검증", "별도 작업", "동일 도구 재사용"),
        ("자산 누적", "0", "catalog + capability + prompt"),
    ]
    table_top = Inches(2.0)
    row_h = Inches(0.6)
    cols = [Inches(3.0), Inches(4.7), Inches(4.7)]
    total_w = sum(cols, Inches(0))
    start_x = (SW - total_w) / 2

    for ri, row in enumerate(rows):
        x = start_x
        is_header = ri == 0
        bg_row = NAVY if is_header else (WHITE if ri % 2 == 1 else SOFT_BLUE)
        text_color = WHITE if is_header else NAVY
        y = table_top + row_h * ri
        for ci, (w, val) in enumerate(zip(cols, row)):
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, row_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = bg_row
            rect.line.color.rgb = GRAY
            tf = rect.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = val
            p.font.size = Pt(13)
            p.font.bold = is_header or ci == 0
            p.font.color.rgb = text_color
            p.alignment = PP_ALIGN.LEFT
            tf.margin_left = Inches(0.15)
            x += w

    add_text_box(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.5),
                 "✓ Settings 신규 앱 진입 → SMOKE_01 PASS까지 약 1시간 (오늘 실측)",
                 size=16, bold=True, color=GREEN)
    add_text_box(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
                 "✓ 누적되는 자산: prompt / capability / catalog / 도구(audit·delta)",
                 size=16, bold=True, color=GREEN)


# ── Slide 11 — 한계와 다음 단계 ─────────────────────────
def slide_next():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title_bar(s, "한계와 다음 단계 — 솔직하게")

    add_card(s, Inches(0.7), Inches(1.4), Inches(6.0), Inches(2.6),
             "현재 한계", [
                 "실 단말·SIM·통신사 lab은 여전히 사람이 셋업",
                 "OS·제조사·통신사 customization은 anchor에 영향",
                 "복잡한 동적 화면(영상·스트림)은 catalog 활용 제한",
                 "AI가 만든 TC는 사람의 “승인 게이트”가 필수",
             ], title_color=RED, bg=SOFT_RED,
             title_size=20, body_size=14)
    add_card(s, Inches(7.0), Inches(1.4), Inches(5.6), Inches(2.6),
             "현재 적용 안 한 것", [
                 "policy v2 운영 자동화 (proposal 단계)",
                 "Tier 0 자동 게이트 (수동 게이트 유지)",
                 "synthetic delta runtime gate",
                 "CI / pre-commit 통합",
             ], title_color=GRAY, bg=LIGHT_BG,
             title_size=20, body_size=14)
    add_card(s, Inches(0.7), Inches(4.2), Inches(6.0), Inches(2.6),
             "단기 다음 단계", [
                 "PR 7B — Markdown report + 추가 fixture",
                 "PR 8 — anchor recommender (오류 anchor 사전 차단)",
                 "Settings SMOKE_02+ — 검색·하위 메뉴",
                 "단말 횡 비교 (AT-M140 / SM-A235N)",
             ], title_color=BLUE, bg=SOFT_BLUE,
             title_size=20, body_size=14)
    add_card(s, Inches(7.0), Inches(4.2), Inches(5.6), Inches(2.6),
             "사내 확산 가능성", [
                 "다른 팀 앱 import → 1주 내 SMOKE_01 도달",
                 "BUG 트래킹 → 단일 repo·문서 표준",
                 "QA 신입 온보딩 자료 자동 생성",
                 "“데이터가 남는 자동화” 표준화",
             ], title_color=GREEN, bg=SOFT_GREEN,
             title_size=20, body_size=14)


# ── Slide 12 — 마무리 + Q&A ─────────────────────────────
def slide_closing():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, NAVY)
    add_text_box(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.7),
                 "정리 — tc-runner가 던지는 3가지 메시지",
                 size=28, bold=True, color=WHITE)

    msgs = [
        ("①", "AI는 “계획자 + 실행자”로 나눠 쓸 때 강하다",
         "Opus가 큰 그림, Sonnet이 디테일 — 컨텍스트가 휘발하지 않는다."),
        ("②", "데이터가 남지 않는 자동화는 거부한다",
         "catalog · reports · prompt = 다음 작업의 입력. 휘발 자동화는 학습이 안 쌓인다."),
        ("③", "1번째 앱이 비싸도, 2번째부터 빠르다",
         "Settings 신규 도입 → SMOKE_01 PASS 약 1시간. 이게 누적의 힘."),
    ]
    y = Inches(1.9)
    for num, head, body in msgs:
        add_text_box(s, Inches(0.8), y, Inches(0.8), Inches(0.7),
                     num, size=44, bold=True, color=ORANGE)
        add_text_box(s, Inches(1.6), y + Inches(0.05),
                     Inches(11.0), Inches(0.6),
                     head, size=22, bold=True, color=WHITE)
        add_text_box(s, Inches(1.6), y + Inches(0.6),
                     Inches(11.0), Inches(0.5),
                     body, size=15, color=SOFT_BLUE)
        y += Inches(1.5)

    add_text_box(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6),
                 "Q & A",
                 size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Run all slides ──────────────────────────────────────
def main():
    slide_title()
    slide_problem()
    slide_overview()
    slide_ai_agents()
    slide_ai_loop()
    slide_pipeline()
    slide_case_music()
    slide_case_others()
    slide_cost_one()
    slide_cost_two()
    slide_next()
    slide_closing()
    prs.save(OUT)
    print(f"✓ saved: {OUT}")
    print(f"  slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
