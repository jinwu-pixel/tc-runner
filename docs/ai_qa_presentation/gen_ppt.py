# -*- coding: utf-8 -*-
"""
사내 AI 활용 발표용 PPTX 생성기 (15장, 16:9, 발표자 노트 포함).

스파인: "AI는 QA 파이프라인의 양 끝에서 일한다."
core 10장 + 증거 5장(E1~E5): json setting TC / MMI focusrule TC /
자동화 PASS·FAIL / BUG-25796 실측표·로그·단말이미지 / BUG-25175 before-after.

재생성: venv\\Scripts\\python.exe docs\\ai_qa_presentation\\gen_ppt.py
S5 성공률은 의도적으로 ▣/▣ placeholder (최신값 확정 전 임의 숫자 금지).
출처 데이터: doc\\BUG25796...\\*.xlsx, output\\BUG25175_Test_Result.xlsx,
reports\\App_Feature_TC_Report.xlsx, tc_samples\\folder_basic_nav.yaml,
ODIN2 - Settings\\SETTINGS_SMOKE_01_app_launch.yaml, menu_tree_settings_*.json
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = r"c:\Users\momen\Projects\tc-runner"
IMG_DEVICE = os.path.join(REPO, "doc", "BUG25796_Mobile_Data_sync_issue", "11.png")

# ---- palette ----
NAVY   = RGBColor(0x1F, 0x2A, 0x44)
AI     = RGBColor(0x1C, 0x7E, 0xD6)
AUTO   = RGBColor(0x6B, 0x72, 0x80)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xB8, 0x6A, 0x00)
TEXT   = RGBColor(0x22, 0x26, 0x2B)
MUTED  = RGBColor(0x5A, 0x63, 0x70)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xEE, 0xF2, 0xF8)
CODEBG = RGBColor(0xF5, 0xF6, 0xF8)
CODEBD = RGBColor(0xCF, 0xD6, 0xDF)
FONT   = "맑은 고딕"
MONO   = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size=18, bold=False, color=TEXT, level=0,
         align=PP_ALIGN.LEFT, space_after=6, first=False, name=FONT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = name
    return p


def title(s, text, eyebrow=None):
    if eyebrow:
        tf = box(s, 0.6, 0.30, 12.1, 0.4)
        para(tf, eyebrow, size=12, bold=True, color=AI, first=True)
    tf = box(s, 0.6, 0.62 if eyebrow else 0.5, 12.13, 0.9)
    para(tf, text, size=26, bold=True, color=NAVY, first=True)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.46),
                              Inches(2.2), Pt(3))
    rule.fill.solid(); rule.fill.fore_color.rgb = AI; rule.line.fill.background()


def bullets(s, items, top=1.75, left=0.75, width=11.9, height=4.4):
    tf = box(s, left, top, width, height)
    for i, (text, lvl) in enumerate(items):
        mark = "•  " if lvl == 0 else "–  "
        para(tf, mark + text, size=(19 if lvl == 0 else 16),
             bold=(lvl == 0), color=(TEXT if lvl == 0 else MUTED),
             level=lvl, space_after=(9 if lvl == 0 else 4), first=(i == 0))
    return tf


def footer(s, text, color=AI, top=6.5, height=0.7):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6),
                             Inches(top), Inches(12.13), Inches(height))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT


def block(s, l, t, w, h, head, sub, fill):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                             Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.color.rgb = fill
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = head
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(10.5); r2.font.color.rgb = WHITE; r2.font.name = FONT
    return shp


def arrow(s, l, t, w, h, color=MUTED):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t),
                           Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background()


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def codeblock(s, l, t, w, h, lines, header=None):
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                              Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = CODEBG
    rect.line.color.rgb = CODEBD
    tf = box(s, l + 0.18, t + 0.12, w - 0.36, h - 0.2)
    first = True
    if header:
        para(tf, header, size=11, bold=True, color=AI, first=True, space_after=4)
        first = False
    for ln in lines:
        para(tf, ln, size=10.5, color=TEXT, first=first, name=MONO, space_after=2)
        first = False
    return rect


def rcolor(text):
    t = text.upper()
    if "FAIL" in t:
        return RED
    if "PASS" in t or "정상" in text:
        return GREEN
    if "mitig" in text.lower():
        return ORANGE
    return TEXT


def cell(tbl, r, c, text, size=11, bold=False, color=TEXT, fill=WHITE,
         align=PP_ALIGN.CENTER):
    cl = tbl.cell(r, c)
    cl.fill.solid(); cl.fill.fore_color.rgb = fill
    cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    cl.margin_top = Pt(2); cl.margin_bottom = Pt(2)
    p = cl.text_frame.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = FONT


def table(s, l, t, headers, rows, widths, top_h=0.42, row_h=0.42,
          result_cols=()):
    nr, nc = len(rows) + 1, len(headers)
    total_w = sum(widths)
    gt = s.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(total_w),
                            Inches(top_h + row_h * len(rows))).table
    gt.first_row = False; gt.horz_banding = False
    for c, w in enumerate(widths):
        gt.columns[c].width = Inches(w)
    for c, h in enumerate(headers):
        cell(gt, 0, c, h, size=11, bold=True, color=WHITE, fill=NAVY)
    for ri, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            if c in result_cols:
                cell(gt, ri, c, val, bold=True, color=rcolor(val))
            else:
                cell(gt, ri, c, val, color=TEXT,
                     align=(PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER))
    return gt


# =========================================================================
# S1 · 표지
# =========================================================================
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.55),
                          Inches(13.333), Inches(2.4))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
para(box(s, 0.8, 2.7, 11.7, 1.0), "AI로 QA를 한다는 것", size=44, bold=True,
     color=WHITE, align=PP_ALIGN.CENTER, first=True)
para(box(s, 0.8, 3.85, 11.7, 0.7), "입력 해석  +  자동화 실행  +  증거 판단",
     size=20, color=LIGHT, align=PP_ALIGN.CENTER, first=True)
para(box(s, 0.8, 5.25, 11.7, 0.7), "AI는 QA 파이프라인의 양 끝에서 일한다",
     size=18, bold=True, color=AI, align=PP_ALIGN.CENTER, first=True)
para(box(s, 0.8, 6.6, 11.7, 0.5), "사내 AI 활용 발표  ·  2026-06-04",
     size=12, color=MUTED, align=PP_ALIGN.CENTER, first=True)
notes(s, "오늘은 'AI로 테스트를 자동화했다'는 단순한 이야기를 하려는 게 아닙니다. "
         "AI를 QA 파이프라인의 어디에 쓰면 실제 가치가 나오는지를 두 지점으로 나눠 "
         "보여드리고, 실제 Excel PASS/FAIL 결과와 TC·로그·단말 화면까지 근거로 "
         "함께 보여드리겠습니다.")

# =========================================================================
# S2 · 한 장 요약
# =========================================================================
s = slide()
title(s, "한 장 요약 — AI는 양 끝에 있다", eyebrow="S2 · OVERVIEW")
by, bh, bw = 1.95, 1.7, 3.3
block(s, 1.16, by, bw, bh, "입력단 AI", "비정형 자료 → TC 해석", AI)
arrow(s, 4.46, by + 0.55, 0.55, 0.6)
block(s, 5.01, by, bw, bh, "Appium 실행", "정형 TC 반복 실행", AUTO)
arrow(s, 8.31, by + 0.55, 0.55, 0.6)
block(s, 8.86, by, bw, bh, "출력단 AI", "증거 → 오판정 차단", AI)
bullets(s, [
    ("입력단: 비정형 자료(PDF/Figma/MMI)를 실행 가능한 TC로 해석", 0),
    ("중간: 정형 TC를 넓게 반복 실행", 0),
    ("출력단: 로그·증거를 해석해 잘못된 결론을 차단", 0),
], top=4.1, height=2.0)
footer(s, "정해진 건 넓게, 정해지지 않은 건 깊게")
notes(s, "결론을 먼저 말씀드립니다. AI의 가치는 가운데 자동화 실행이 아니라 양 끝에 "
         "있습니다. 들어올 때 비정형 자료를 테스트로 바꾸고, 나갈 때 증거를 읽어 "
         "판단합니다. 자동화는 정해진 걸 넓게, AI는 정해지지 않은 걸 깊게.")

# =========================================================================
# S3 · 입력의 진화와 PDF의 한계
# =========================================================================
s = slide()
title(s, "입력의 진화와 PDF의 한계", eyebrow="S3 · 입력단 AI")
para(box(s, 0.7, 1.62, 5.6, 0.35), "Before — 칸 밀림", size=14, bold=True,
     color=RED, first=True)
tb = s.shapes.add_table(3, 3, Inches(0.7), Inches(2.0), Inches(5.6),
                        Inches(1.5)).table
for c, h in enumerate(["No.", "분류", "값"]):
    cell(tb, 0, c, h, size=12, bold=True, color=WHITE, fill=AUTO)
cell(tb, 1, 0, "14"); cell(tb, 1, 1, "버튼"); cell(tb, 1, 2, "확인")
cell(tb, 2, 0, "15")
cell(tb, 2, 1, "취소", bold=True, color=WHITE, fill=RED)
cell(tb, 2, 2, "999", bold=True, color=WHITE, fill=RED)
para(box(s, 0.7, 3.6, 5.6, 0.5), "한 줄 = 3칸 가정 → 팝업 줄(설명 2칸) 뒤로 한 칸씩 밀림",
     size=12, color=MUTED, first=True)
para(box(s, 7.0, 1.62, 5.6, 0.35), "After — No. 기준 복구", size=14, bold=True,
     color=GREEN, first=True)
tb = s.shapes.add_table(3, 3, Inches(7.0), Inches(2.0), Inches(5.6),
                        Inches(1.5)).table
for c, h in enumerate(["No.", "분류", "값"]):
    cell(tb, 0, c, h, size=12, bold=True, color=WHITE, fill=AUTO)
cell(tb, 1, 0, "14"); cell(tb, 1, 1, "버튼"); cell(tb, 1, 2, "확인")
cell(tb, 2, 0, "15")
cell(tb, 2, 1, "버튼", bold=True, color=WHITE, fill=GREEN)
cell(tb, 2, 2, "취소", bold=True, color=WHITE, fill=GREEN)
para(box(s, 7.0, 3.6, 5.6, 0.5), "No.(번호) 칸을 기준점으로 줄 분리 → 제자리 복구",
     size=12, color=MUTED, first=True)
bullets(s, [
    ("입력의 진화: PDF → Figma JSON → MMI/Excel", 0),
    ("PDF는 사람이 보기 좋지만 기계가 읽기 어렵다", 0),
    ("결과: 추출 정상화 10건 복구 (stale → 정상)", 0),
], top=4.35, height=2.0)
notes(s, "모든 자동화는 '무엇을 테스트할지'라는 입력에서 시작합니다. 실제로 표를 줄 "
         "단위로 잘랐더니, 설명이 두 칸인 팝업 줄 때문에 그 뒤가 한 칸씩 밀려서 취소 "
         "버튼 글씨와 숫자가 분류값으로 들어갔습니다. 번호 칸을 기준점으로 바꿔 "
         "해결했고, 비어 있던 추출 결과 10건이 정상으로 살아났습니다.")

# =========================================================================
# S4 · Figma JSON / MMI 정규화
# =========================================================================
s = slide()
title(s, "Figma JSON / MMI 정규화", eyebrow="S4 · 입력단 AI")
bullets(s, [
    ("Figma JSON: frame/component/text 계층, structural anchor, expected text 추출 용이", 0),
    ("MMI/Excel: 기존 TC 자산 재사용 가능", 0),
    ("함정: design source라 실제 device와 다를 수 있음", 0),
    ("locale · variant · hidden state · runtime permission 은 별도 검증", 1),
    ("해결: 소스별 adapter → 공통 TCIR / ExtractionRecord 로 정규화", 0),
], top=1.8, height=4.0)
footer(s, "Figma JSON은 'PDF보다 똑똑한 입력'이지 '정답'은 아니다")
notes(s, "Figma JSON으로 넘어가면 화면 구조·컴포넌트·텍스트를 기계가 훨씬 잘 읽습니다. "
         "MMI 같은 기존 자산도 변환 규칙만 맞추면 재사용됩니다. 다만 Figma는 디자인 "
         "원본이라 실제 단말과 다를 수 있어서 언어·변형·숨은 상태는 따로 검증해야 "
         "합니다. 그래서 소스가 무엇이든 공통 형식으로 정규화합니다.")

# =========================================================================
# E1 · 입력 예시 ① — JSON 메뉴트리 → Settings TC
# =========================================================================
s = slide()
title(s, "입력 예시 ① — JSON 메뉴트리 → Settings TC", eyebrow="E1 · 입력단 예시")
codeblock(s, 0.7, 1.75, 5.7, 3.7, [
    '"focus": "...settings.Settings",',
    '"fingerprint": "d15a7f0e",',
    '"nodes": [',
    '  { "text": "네트워크 및 인터넷",',
    '    "bounds": "[0,176][480,256]" },',
    '  { "text": "연결된 기기" },',
    '  { "text": "앱" },',
    '  { "text": "알림" },',
    '  { "text": "배터리" } ]',
], header="menu_tree_settings_*.json  (device-sourced 추출)")
arrow(s, 6.5, 3.3, 0.45, 0.6, color=AI)
codeblock(s, 7.05, 1.75, 5.55, 3.7, [
    'tc_name: SETTINGS_SMOKE_01_app_launch',
    '- action: shell',
    '    command: am start ...settings/.Settings',
    '- action: verify_text  target: "설정"',
    '- action: verify_text  target: "네트워크 및 인터넷"',
    '- action: verify_text  target: "연결된 기기"',
    '- action: verify_text  target: "앱"',
    '- action: verify_text  target: "알림"',
    '- action: verify_text  target: "배터리"',
], header="생성된 TC  (com.android.settings ROOT 라벨 정합)")
footer(s, "device-sourced 메뉴트리 JSON = figma ExtractionRecord와 평행 구조 → TC 합성 입력",
       top=5.7, height=0.65)
para(box(s, 0.7, 6.45, 12.0, 0.5),
     "사용자 단(설정 ROOT)에서 보이는 top-level 라벨을 JSON에서 추출 → 그대로 verify_text TC로",
     size=12, color=MUTED, first=True)
notes(s, "입력단 AI의 첫 실물 예시입니다. 왼쪽은 실기에서 read-only로 추출한 설정 앱 "
         "메뉴트리 JSON입니다. focus, fingerprint, 그리고 노드별 텍스트와 좌표가 들어 "
         "있죠. 이 JSON은 Figma의 ExtractionRecord와 평행 구조라, 오른쪽처럼 설정 "
         "ROOT의 top-level 라벨 — 네트워크, 연결된 기기, 앱, 알림, 배터리 — 을 그대로 "
         "verify_text TC로 합성합니다. 즉 사용자가 화면에서 보는 것을 입력 삼아 "
         "테스트가 자동으로 만들어집니다.")

# =========================================================================
# E2 · 입력 예시 ② — MMI focusrule TC (폴더폰 방향키)
# =========================================================================
s = slide()
title(s, "입력 예시 ② — MMI focusrule TC (폴더폰 방향키)", eyebrow="E2 · 입력단 예시")
codeblock(s, 0.7, 1.75, 6.3, 3.6, [
    'tc_name: FOLDER_BASIC_NAV_01',
    '- action: key            key: 3        # HOME',
    '- action: wait           duration: 1000',
    '- action: verify_focus_moved',
    '    trigger_action: key',
    '    trigger_step: { key: 20 }   # DPAD_DOWN',
    '- action: key_sequence   keys: [20, 22]  # ↓ →',
    '    delay: 0.5',
    '- action: key            key: 3',
], header="tc_samples/folder_basic_nav.yaml")
bullets(s, [
    ("폴더폰/피처폰은 터치가 아니라 방향키(MMI)로 이동", 0),
    ("verify_focus_moved = '방향키 입력 후 포커스가 실제로 이동했는지' 검증", 0),
    ("key_sequence 로 다중 키 깊이 진입(↓↓→ 등)", 0),
    ("위험 팝업은 기본 포커스가 'Cancel'인지까지 검증(F4 패턴)", 0),
], top=1.85, left=7.2, width=5.5, height=3.6)
footer(s, "사용자 단(피처폰)에서 보이는 방향키 동작을 그대로 TC로", top=5.6, height=0.65)
para(box(s, 0.7, 6.4, 12.0, 0.5),
     "schema: verify_focus_moved · key_sequence (strict focus-moved / key sequence actions)",
     size=12, color=MUTED, first=True)
notes(s, "두 번째 입력 예시는 폴더폰·피처폰입니다. 이런 단말은 터치가 아니라 방향키로 "
         "움직이죠. 그래서 verify_focus_moved 라는 액션으로 '방향키를 눌렀을 때 포커스가 "
         "실제로 이동했는가'를 검증합니다. key_sequence로 여러 키를 연속 입력해 깊은 "
         "메뉴로 들어가고, 삭제 같은 위험 팝업에선 기본 포커스가 Cancel에 있는지까지 "
         "확인합니다. 사용자가 손으로 누르는 동작을 그대로 TC로 옮긴 것입니다.")

# =========================================================================
# S5 · 중간층 — Appium 자동화 실행
# =========================================================================
s = slide()
title(s, "중간층 — Appium 자동화 실행", eyebrow="S5 · 자동화 실행")
bullets(s, [
    ("정형 TC → 실행 가능한 Appium 테스트로 변환", 0),
    ("CI/CD · 회귀 반복 · 리포트 · Jira 연동", 0),
    ("강점: 양과 반복", 0),
    ("지표: 성공률  ▣ / ▣   ← 발표 직전 최신값 확정 후 삽입", 0),
], top=1.85, height=4.0)
footer(s, "자동화는 정해진 것을 많이, 빠르게", color=AUTO)
notes(s, "입력이 정리되면, 정형화된 부분은 Appium이 실제 실행 가능한 테스트로 바꿔 "
         "반복 실행합니다. CI/CD와 리포트, Jira에 붙어 회귀를 넓게 커버합니다. "
         "성공률 수치는 발표 직전에 최신값으로 채워 넣겠습니다.")

# =========================================================================
# E3 · 자동화 실행 결과 — TC PASS / FAIL
# =========================================================================
s = slide()
title(s, "자동화 실행 결과 — TC PASS / FAIL", eyebrow="E3 · 자동화 결과")
table(s, 0.7, 1.8,
      ["구분", "건수", "결과"],
      [["기능 TC (Gallery 14 + Memosy 14)", "28", "28 PASS"],
       ["하드키 TC", "18", "17 PASS · 1 FAIL"]],
      [5.2, 1.4, 3.0], result_cols=(2,))
para(box(s, 0.7, 3.35, 9.6, 0.35), "대표 케이스", size=13, bold=True, color=NAVY,
     first=True)
table(s, 0.7, 3.75,
      ["TC", "내용", "판정"],
      [["G-01", "사진 목록(날짜별 그룹핑) 표시", "PASS"],
       ["HK-G05", "ENTER → 사진 상세보기 진입", "PASS"],
       ["HK-G12", "하단 탭 D-pad 포커스 도달", "FAIL"]],
      [1.4, 6.8, 1.4], result_cols=(2,))
para(box(s, 0.7, 6.4, 12.0, 0.5),
     "출처: reports/App_Feature_TC_Report.xlsx · AT-M140 · ADB+UIAutomator 자동화",
     size=12, color=MUTED, first=True)
notes(s, "자동화가 실제로 낸 결과입니다. 기능 TC는 갤러리와 메모지 합쳐 28건 전부 "
         "PASS, 하드키 TC는 18건 중 17건 PASS에 1건 FAIL입니다. FAIL인 HK-G12는 하단 "
         "탭에 D-pad 포커스가 도달하지 못하는 실제 결함이고, 자동화가 이를 그냥 넘기지 "
         "않고 FAIL로 잡아낸 것이 핵심입니다.")

# =========================================================================
# S6 · 자동화가 멈춰야 하는 순간
# =========================================================================
s = slide()
title(s, "자동화가 멈춰야 하는 순간", eyebrow="S6 · 자동화 실행")
bullets(s, [
    ("입력이 모호하면 자동화도 모호해진다", 0),
    ("멈춰야 하는 신호:", 0),
    ("언어 불일치 — 단말 한국어 ↔ 기대값 일본어 (FR-HK-070)", 1),
    ("파괴적 동작 (destructive action)", 1),
    ("모호한 expected result", 1),
], top=1.85, height=4.0)
footer(s, "좋은 자동화는 '무조건 누르기'가 아니라 멈출 줄 안다")
notes(s, "자동화는 입력이 모호하면 같이 모호해집니다. 단말 언어가 한국어인데 기대값이 "
         "일본어인 케이스는 기계적으로 통과시키면 안 됩니다. 좋은 자동화의 조건은 빠른 "
         "실행이 아니라 멈출 줄 아는 것 — 여기서 출력단 AI가 필요해집니다.")

# =========================================================================
# S7 · 출력단 AI — tc-runner는 학습 루프
# =========================================================================
s = slide()
title(s, "출력단 AI — tc-runner는 학습 루프", eyebrow="S7 · 출력단 AI")
steps = ["탐색", "catalog", "TC", "runtime 증거", "delta", "다음 판단"]
lx, top, bw, gap, h = 0.85, 1.95, 1.55, 0.35, 0.95
for i, st in enumerate(steps):
    l = lx + i * (bw + gap)
    block(s, l, top, bw, h, st, None, AI if i in (0, 3, 5) else AUTO)
    if i < len(steps) - 1:
        arrow(s, l + bw - 0.02, top + 0.30, gap + 0.04, 0.36)
bullets(s, [
    ("단순 실행기가 아니다", 0),
    ("모든 실행이 데이터로 남는다", 0),
    ("다음 테스트 설계가 빨라진다 (복리 효과)", 0),
    ("AI의 가치는 '누르는 것'보다 '판단을 누적하는 것'", 0),
], top=3.4, height=2.9)
notes(s, "출력단의 tc-runner는 한 번 돌리고 끝나는 실행기가 아니라 학습 루프입니다. "
         "탐색한 화면·셀렉터·실패 원인을 카탈로그로 쌓고, 다음 빌드에선 전수 재탐색 "
         "대신 바뀐 부분만 봅니다. 매번 데이터가 남아 다음 판단이 빨라집니다.")

# =========================================================================
# S8 · 증거 장면 ① — 잘못된 전제 교정 (BUG-17126)
# =========================================================================
s = slide()
title(s, "증거 장면 ① — 잘못된 전제 교정 (BUG-17126)", eyebrow="S8 · 출력단 AI")
bullets(s, [
    ("SM Reject cause 27 = PS attach 성공 후 잘못된 APN으로 PDP만 reject 되는 구조", 0),
    ("KT 미인증 SIM 조건에서는 목표 상태(SM reject cause 27)에 도달하지 못함", 0),
    ("판단: 인증 SKT USIM 필수로 테스트 설계를 교정 — 단말을 연결하기도 전에", 0),
], top=1.85, height=3.6)
footer(s, "AI가 테스트 실행 전에 테스트 전제를 고쳤다")
notes(s, "SM Reject cause 27은 접속 성공 뒤 PDP만 거부되는 구조입니다. 기존 가정은 "
         "KT 미인증 SIM이었는데, 그 조건에서는 이 상태에 도달하지 못합니다. AI가 "
         "프로토콜 구조를 추론해 단말을 꽂기도 전에 '인증 SKT USIM이라야 성립한다'고 "
         "전제를 바로잡았습니다.")

# =========================================================================
# S9 · 증거 장면 ② — 단말 단독 결함이 아님을 입증 (BUG-25796)
# =========================================================================
s = slide()
title(s, "증거 장면 ② — 단말 단독 결함이 아님을 입증 (BUG-25796)",
      eyebrow="S9 · 출력단 AI")
bullets(s, [
    ("34사이클 매트릭스", 0),
    ("레퍼런스 단말(SM-A235N) 대조", 0),
    ("WWAN 정/역 재현: on trigger ↔ off 0", 0),
    ("정량: 130.66s  σ=0.13s", 0),
], top=1.85, height=2.7)
para(box(s, 0.75, 4.55, 11.9, 0.5),
     "작은 주석: 신빌드에서 modem recovery는 개선됐지만 host trigger 자체는 잔존",
     size=12, color=MUTED, first=True)
footer(s, "결론: 단말 단독 결함이 아니라, 호스트 Windows WWAN AutoConfig 간섭으로 규명")
notes(s, "DataPopup 이슈는 처음엔 단말 결함처럼 보였습니다. 34사이클 매트릭스, 정상 "
         "단말 대조, WWAN 양방향 재현, 130.66초 σ 0.13초 정량까지 거쳐, 단말 단독 "
         "결함이 아니라 호스트 윈도우 WWAN 간섭으로 규명했습니다. 다음 장에서 실측 "
         "데이터를 보여드립니다.")

# =========================================================================
# E4 · BUG-25796 실측 — 매트릭스 · 로그 · 단말 화면
# =========================================================================
s = slide()
title(s, "BUG-25796 실측 — WWAN 매트릭스 · 로그 · 단말 화면",
      eyebrow="E4 · 출력단 증거")
table(s, 0.7, 1.7,
      ["빌드", "WWAN", "tap→성공", "0xffff", "판정"],
      [["Z0527U (신)", "OFF", "~1.4s", "0", "정상"],
       ["Z0527U (신)", "ON", "~4s", "6 (transient)", "mitigated"],
       ["Z0518U (구)", "ON", "~125s", "48 (loop)", "sustained FAIL"],
       ["REF A235N", "ON", "~74s", "0", "정상"],
       ["REF A235N", "OFF", "~70s", "0", "정상"]],
      [1.9, 1.0, 1.3, 1.9, 2.2], result_cols=(4,))
codeblock(s, 0.7, 4.55, 9.3, 2.3, [
    'WWAN-OFF · 6/1 신빌드   12:21:16 tap → 12:21:18 PDN  0xffff=0   ~1.4s 정상',
    'WWAN-ON  · 6/1 신빌드   14:27:47 tap → 14:27:51 PDN  0xffff=6   transient → 회복',
    'WWAN-ON  · 구빌드 Z0518U 16:01:38 tap → 16:03:42 PDN  0xffff=48  ~125s sustained FAIL',
], header="AP logcat 타임라인 (RILJ SETUP_DATA_CALL · 0xffff=ERROR_UNSPECIFIED)")
if os.path.exists(IMG_DEVICE):
    s.shapes.add_picture(IMG_DEVICE, Inches(10.4), Inches(2.55), height=Inches(4.3))
    para(box(s, 10.4, 1.72, 2.5, 0.6),
         "ODIN2 / LG U+\n모바일 데이터 켜짐", size=10, color=MUTED, first=True)
notes(s, "이게 실측 데이터입니다. 표를 보면 신빌드는 WWAN OFF에서 1.4초, ON에서도 "
         "transient 6건 뒤 즉시 회복됩니다. 반면 구빌드 Z0518U는 125초에 0xffff가 48번 "
         "반복되는 sustained FAIL이고, 정상 레퍼런스 단말은 ON/OFF 모두 0으로 깨끗합니다. "
         "아래는 실제 logcat 타임라인, 오른쪽은 그때의 ODIN2 단말 화면입니다. 즉 "
         "사용자 체감 버그는 mitigated 됐지만 host trigger 자체는 남아 있다는 결론입니다.")

# =========================================================================
# E5 · BUG-25175 회귀 — 수정 전 FAIL → 수정 후 PASS
# =========================================================================
s = slide()
title(s, "BUG-25175 회귀 — 수정 전 FAIL → 수정 후 PASS",
      eyebrow="E5 · 출력단 증거")
table(s, 0.7, 1.8,
      ["#", "요구사항", "수정 전", "수정 후"],
      [["1", "APN 진입 시 Internet PDN만 표시", "FAIL", "PASS"],
       ["2", "APN 추가 불가 (+ 버튼 제거)", "FAIL", "PASS"],
       ["3", "Internet PDN APN 수정/삭제 불가", "FAIL", "PASS"],
       ["4", "IMS·Tethering PDN 등 비노출", "FAIL", "PASS"],
       ["5", "엔지니어 모드 전체 APN 표시", "N/A", "PASS"],
       ["6", "엔지니어 모드 APN 추가/수정", "N/A", "PASS"],
       ["7", "타사 USIM(SKT/KT) 기존 메뉴 유지", "PASS", "PASS"]],
      [0.7, 7.0, 1.7, 1.7], result_cols=(2, 3))
footer(s, "전체 결과: 18/18 ALL PASS  ·  3-way ground truth(UI · dumpsys · 인터페이스) 일치",
       top=6.45, height=0.65, color=GREEN)
notes(s, "BUG-25175는 LGU+ APN 메뉴 회귀입니다. 수정 전에는 요구사항 1~4가 모두 FAIL "
         "이었는데, 수정 후 전부 PASS로 전환됐습니다. 엔지니어 모드 항목은 이전엔 메뉴 "
         "자체가 없어 N/A였다가 이제 PASS입니다. 전체 18개 항목이 ALL PASS이고, 단말 "
         "UI·dumpsys·인터페이스 3중으로 교차 확인했습니다.")

# =========================================================================
# S10 · 신뢰 장치와 결론
# =========================================================================
s = slide()
title(s, "신뢰 장치와 결론", eyebrow="S10 · 마무리")
guards = [
    ("PASS 어휘 구분", "validate / runtime / manual evidence / BUG-GAP"),
    ("read-only guard", "읽기전용 설계 · 정적검사 · 런타임검증"),
    ("provenance / confidence", "source 출처 · 신뢰도 · manual-required reason"),
    ("scope NOTE", "외부 요인을 FAIL과 격리"),
]
gx, gy, gw, gh = 0.75, 1.85, 5.85, 1.05
for i, (head, sub) in enumerate(guards):
    row, col = divmod(i, 2)
    l = gx + col * (gw + 0.35)
    t = gy + row * (gh + 0.3)
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                             Inches(gw), Inches(gh))
    shp.fill.solid(); shp.fill.fore_color.rgb = LIGHT; shp.line.color.rgb = AI
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = head
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sub
    r2.font.size = Pt(11); r2.font.color.rgb = MUTED; r2.font.name = FONT
concl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.35),
                           Inches(12.13), Inches(1.0))
concl.fill.solid(); concl.fill.fore_color.rgb = NAVY; concl.line.fill.background()
tf = concl.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "AI는 테스트를 대충 통과시키는 도구가 아니라, 입력을 해석하고 증거를 판정해 가짜 통과와 오판정을 막는 장치다"
r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
para(box(s, 0.6, 6.55, 12.13, 0.6), "자동화는 넓게, AI는 깊게.", size=18, bold=True,
     color=AI, align=PP_ALIGN.CENTER, first=True)
notes(s, "AI를 신뢰하게 만드는 건 모델 자체가 아니라 둘레의 장치들입니다. 통과를 네 "
         "가지로 구분하고, AI가 실수로도 폰을 못 건드리게 잠그고, 결과마다 출처와 "
         "신뢰도를 남기고, 외부 요인은 FAIL과 격리합니다. 안내+관찰 도구는 자동화가 "
         "부적절하면 대신 누르지 않고, 언어 불일치를 가짜 통과로 넘기지 않았습니다. "
         "한 문장으로 맺습니다 — AI는 입력을 해석하고 증거를 판정해 가짜 통과와 "
         "오판정을 막아줍니다. 자동화는 넓게, AI는 깊게.")

# ---- save (잠금 시 새 파일명으로 폴백) ----
BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "AI_QA_발표_2026-06-04.pptx")
try:
    prs.save(OUT)
except PermissionError:
    n = 2
    while True:
        alt = os.path.join(BASE, f"AI_QA_발표_2026-06-04_v{n}.pptx")
        try:
            prs.save(alt)
            OUT = alt
            print("LOCKED primary (열려 있음?) -> 새 파일로 저장")
            break
        except PermissionError:
            n += 1
print("SAVED:", OUT)
print("SLIDES:", len(prs.slides._sldIdLst))
