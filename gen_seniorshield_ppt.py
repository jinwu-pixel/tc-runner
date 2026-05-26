"""SeniorShield 사내 배포용 PPT 생성 스크립트. (rev2 — 6건 필수 + 1건 권장 반영)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 설정 ──────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "doc", "SeniorShield_소개자료.pptx")

PHONE_W = Inches(2.2)
PHONE_H = Inches(3.67)
PHONE_H_SM = Inches(3.0)  # 캡션 겹침 방지용 축소 높이

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
BLUE = RGBColor(0x29, 0x80, 0xB9)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xFA)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
PURPLE_BG = RGBColor(0xE8, 0xDA, 0xEF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, text, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                   prs.slide_width, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)


def add_text_box(slide, left, top, width, height, text, size=18, bold=False,
                 color=NAVY, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, size=16, color=NAVY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def add_phone_img(slide, img_path, left, top, width=PHONE_W, height=PHONE_H):
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width, height)


def add_box(slide, left, top, width, height, fill_color, text="",
            font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER


# ━━━━━━━━━━━ SLIDE 1: 표지 ━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "시니어쉴드 (SeniorShield)", size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
# ★ 권장: "차단" → "사전 감지 및 대응 지원"
add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
             "보이스피싱 사전 감지 및 대응 지원 앱", size=28,
             color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
             "v1.0  |  Android 8.0+  |  사내 배포용", size=20, color=GRAY,
             align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.6),
             "2026-04-17  |  Mobile QA Team", size=16, color=GRAY,
             align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━ SLIDE 2: 앱 개요 ━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "앱 개요")
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(0.6),
             "시니어쉴드는 보이스피싱 피해를 사전에 감지하고 대응을 돕는 Android 앱입니다.",
             size=22, bold=True)
items = [
    "통화, 앱 사용, 앱 설치, 기기 환경을 실시간 모니터링",
    "위험 징후 감지 시 전체화면 경고 팝업으로 즉시 알림",
    "뱅킹앱/텔레뱅킹 이용 시 잠시 멈춤 쿨다운 제공",
    "보이스피싱 대응 연습 시뮬레이션 내장",
    "보호자 연락처 등록 및 긴급 연락 지원",
]
add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(7), Inches(3), items, size=18)
add_box(slide, Inches(0.8), Inches(5.5), Inches(7), Inches(0.7), RED,
        "핵심 원칙: 보호자 자동 알림, SMS 자동 발송 없음. 모든 외부 연락은 사용자가 직접 실행.",
        font_size=15)
add_phone_img(slide, os.path.join(BASE, "c4_after.png"), Inches(9), Inches(1.5))
add_text_box(slide, Inches(9), Inches(5.3), PHONE_W, Inches(0.5),
             "안전 상태 홈 화면", size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━ SLIDE 3: 감지 체계 ━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "감지 체계 — 무엇을 감시하나요?")
sources = [
    ("전화 감시", "미저장 번호 수신\n장시간 통화 (3분+)\n반복 수신\n은행 ARS 발신", BLUE),
    ("앱 사용 감시", "원격제어앱 실행\n(TeamViewer 등)\n원격제어 후\n뱅킹앱 실행", ORANGE),
    ("앱 설치 감시", "사이드로드 앱 설치\n원격제어앱 설치", RED),
    ("기기 환경", "루팅 감지\ntest-keys 감지", GRAY),
]
for i, (title, desc, color) in enumerate(sources):
    left = Inches(0.6 + i * 3.1)
    add_box(slide, left, Inches(1.5), Inches(2.8), Inches(0.6), color, title, 18)
    txBox = slide.shapes.add_textbox(left, Inches(2.3), Inches(2.8), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in desc.split('\n'):
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = "  " + line
        p.font.size = Pt(15)
        p.font.color.rgb = NAVY
        p.space_after = Pt(4)
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(12), Inches(0.6),
             "위험 등급 (점수 기반)", size=22, bold=True)
levels = [
    ("LOW", "< 25점", "미저장 번호 단독", BLUE),
    ("MEDIUM", "25~49점", "텔레뱅킹 단독", ORANGE),
    ("HIGH", "50~79점", "미저장 + 장시간 통화", RGBColor(0xD3, 0x54, 0x00)),
    ("CRITICAL", "80점+", "통화 + 원격제어 조합", RED),
]
for i, (name, score, example, color) in enumerate(levels):
    left = Inches(0.6 + i * 3.1)
    add_box(slide, left, Inches(5.1), Inches(2.8), Inches(0.5), color,
            f"{name}  ({score})", 14)
    add_text_box(slide, left, Inches(5.7), Inches(2.8), Inches(0.5),
                 f"예: {example}", size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━ SLIDE 4: 위험 팝업 ★필수5 통화/비통화 분리 ━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "경고 방식 1 — 위험 팝업", RED)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(0.6),
             "위험 징후 감지 시 전체화면 경고 팝업이 즉시 표시됩니다.", size=20, bold=True)
items = [
    "원격제어앱(TeamViewer 등) 실행 감지 시",
    "의심 통화 후 은행 ARS 전화 발신 시",
    "원격제어 중 뱅킹앱 실행 시",
    "사이드로드 앱 설치 감지 시",
]
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(1.8), items, size=15)

# ★ 필수5: 통화 중 버튼
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.4),
             "팝업 버튼 (통화 중)", size=16, bold=True)
btn_call = [
    "주: 전화 앱으로 이동 — 다이얼러 전환",
    "보조: 통화 경고 닫기 — 같은 통화 중 재경고 억제",
    "등록된 보호자에게 문자 보내기 (보호자 등록 시)",
    "앱 열어서 확인하기 — 시니어쉴드 홈으로",
]
add_bullet_list(slide, Inches(0.8), Inches(4.4), Inches(5.5), Inches(1.5), btn_call, size=13)

# ★ 필수5: 비통화 버튼
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(5.5), Inches(0.4),
             "팝업 버튼 (비통화)", size=16, bold=True)
btn_nocall = [
    "주: 일단 닫기 — 팝업 닫기 (세션 유지)",
    "보조: 위험 경고 해제 — 세션 종료",
    "등록된 보호자에게 문자 보내기 (보호자 등록 시)",
    "앱 열어서 확인하기 — 시니어쉴드 홈으로",
]
add_bullet_list(slide, Inches(0.8), Inches(6.1), Inches(5.5), Inches(1.3), btn_nocall, size=13)

add_phone_img(slide, os.path.join(BASE, "a_popup.png"), Inches(7.5), Inches(1.3))
add_text_box(slide, Inches(7.5), Inches(5.1), PHONE_W, Inches(0.5),
             "HIGH 위험 팝업 (비통화)", size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_phone_img(slide, os.path.join(BASE, "c1.png"), Inches(10.2), Inches(1.3))
add_text_box(slide, Inches(10.2), Inches(5.1), PHONE_W, Inches(0.5),
             "CRITICAL 위험 팝업 (비통화)", size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━ SLIDE 5: 뱅킹 쿨다운 ★필수4 시간 3분류 ━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "경고 방식 2 — 뱅킹 쿨다운", ORANGE)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(6.5), Inches(0.6),
             "의심 상황에서 뱅킹앱 사용 시 카운트다운 화면으로 잠시 멈춤 유도",
             size=20, bold=True)
items = [
    "의심 통화 중 뱅킹앱 실행 → 쿨다운 발동",
    "의심 통화 후 은행 ARS 발신 → 쿨다운 발동",
    "같은 세션에서 쿨다운은 1회만 발동 (반복 방해 방지)",
    "통화 중이면 \"전화 앱으로 이동\" 버튼 제공",
]
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(6.5), Inches(2.0), items, size=16)

# ★ 필수4: 쿨다운 시간 3분류 정합
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(6.5), Inches(0.5),
             "쿨다운 카운트다운 시간", size=18, bold=True)
cd_modes = [
    ("CRITICAL 실세션", "60초", "프로덕션 CRITICAL 위험", RED),
    ("HIGH 실세션", "30초", "프로덕션 HIGH 위험", ORANGE),
    ("디버그 미리보기", "5초", "사내 테스트 전용", GRAY),
]
for i, (mode, sec, desc, color) in enumerate(cd_modes):
    y = Inches(4.8 + i * 0.55)
    add_box(slide, Inches(0.8), y, Inches(2.2), Inches(0.42), color, mode, 13)
    add_text_box(slide, Inches(3.2), y, Inches(1.0), Inches(0.42), sec, size=15, bold=True)
    add_text_box(slide, Inches(4.3), y, Inches(3.0), Inches(0.42), desc, size=13, color=GRAY)

# ★ 필수4: 스크린샷 캡션 정합
add_phone_img(slide, os.path.join(BASE, "b_cool.png"), Inches(7.8), Inches(1.3))
add_text_box(slide, Inches(7.8), Inches(5.1), PHONE_W, Inches(0.5),
             "디버그 미리보기 (5초, 표시값 변동)", size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_phone_img(slide, os.path.join(BASE, "d_cool.png"), Inches(10.5), Inches(1.3))
add_text_box(slide, Inches(10.5), Inches(5.1), PHONE_W, Inches(0.5),
             "실세션 통화 중 쿨다운 (60초 중)", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━ SLIDE 6: 사용 흐름 ★필수6 캡션 겹침 수정 ━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "사용 흐름 — 위험 감지부터 안전 확인까지")
steps = [
    ("1. 위험 감지", "의심 통화/원격제어앱\n실행 감지", RED),
    ("2. 경고 표시", "전체화면 팝업 또는\n쿨다운 카운트다운", ORANGE),
    ("3. 사용자 확인", "경고 확인 후\n안전 여부 판단", BLUE),
    ("4. 안전 확인", "활성 위험 세션 상태에서\n\"안전 확인했어요\" 탭\n→ 정상 복귀", GREEN),
]
for i, (title, desc, color) in enumerate(steps):
    left = Inches(0.5 + i * 3.2)
    add_box(slide, left, Inches(1.5), Inches(2.6), Inches(0.6), color, title, 18)
    add_text_box(slide, left, Inches(2.3), Inches(2.6), Inches(1.2),
                 desc, size=14, color=NAVY, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(slide, left + Inches(2.6), Inches(1.6), Inches(0.6), Inches(0.5),
                     "→", size=28, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# ★ 필수6: 이미지 축소 + 캡션 아래로 분리
img_top = Inches(3.7)
cap_top = img_top + PHONE_H_SM + Inches(0.1)  # 이미지 아래 여백

add_phone_img(slide, os.path.join(BASE, "c3_app_home.png"),
              Inches(1.0), img_top, PHONE_W, PHONE_H_SM)
add_text_box(slide, Inches(1.0), cap_top, PHONE_W, Inches(0.4),
             "위험 감지 상태", size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_phone_img(slide, os.path.join(BASE, "c1.png"),
              Inches(5.5), img_top, PHONE_W, PHONE_H_SM)
add_text_box(slide, Inches(5.5), cap_top, PHONE_W, Inches(0.4),
             "CRITICAL 경고 팝업", size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_phone_img(slide, os.path.join(BASE, "c4_after.png"),
              Inches(10.0), img_top, PHONE_W, PHONE_H_SM)
add_text_box(slide, Inches(10.0), cap_top, PHONE_W, Inches(0.4),
             "안전 확인 후 홈", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━ SLIDE 7: 필수 권한 ━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "설치 후 필수 권한 설정")
permissions = [
    ("전화 / 통화기록 / 연락처 / 발신", "통화 상태 감지, 미저장 번호 판별, 텔레뱅킹 번호 확인",
     "앱 첫 실행 시 다이얼로그"),
    ("알림", "위험 감지 시 즉시 알림", "앱 첫 실행 시 다이얼로그 (Android 13+)"),
    ("사용 정보 접근", "원격제어앱, 뱅킹앱 실행 감지", "설정 > 사용 정보 접근 > 시니어쉴드 ON"),
    ("다른 앱 위에 표시", "위험 팝업, 쿨다운 오버레이 표시", "설정 > 다른 앱 위에 표시 > 시니어쉴드 ON"),
]
add_box(slide, Inches(0.6), Inches(1.5), Inches(3.0), Inches(0.5), NAVY, "권한", 15)
add_box(slide, Inches(3.7), Inches(1.5), Inches(5.0), Inches(0.5), NAVY, "용도", 15)
add_box(slide, Inches(8.8), Inches(1.5), Inches(4.0), Inches(0.5), NAVY, "설정 방법", 15)
for i, (perm, usage, how) in enumerate(permissions):
    y = Inches(2.2 + i * 0.8)
    add_text_box(slide, Inches(0.8), y, Inches(2.8), Inches(0.6), perm, size=15, bold=True)
    add_text_box(slide, Inches(3.9), y, Inches(4.8), Inches(0.6), usage, size=14)
    add_text_box(slide, Inches(9.0), y, Inches(3.8), Inches(0.6), how, size=14, color=BLUE)
add_box(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.7), PURPLE_BG,
        "참고: 사용 정보 접근 권한이 없으면 원격제어앱/뱅킹앱 감지가 불가합니다.",
        font_size=15, font_color=NAVY)

# ━━━━━━━━━━━ SLIDE 8: 디버그 패널 ━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "테스터용 — 디버그 패널 (DEBUG 빌드 전용)")
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(0.5),
             "홈 → 스크롤 하단 → 앱 설정에서 접근", size=18, bold=True)
debug_items = [
    "테스트 모드: 통화 임계값 3분 → 10초로 단축",
    "위험 팝업 미리보기: HIGH 레벨 팝업 강제 표시",
    "뱅킹 쿨다운 미리보기: 5초 쿨다운 강제 표시",
    "텔레뱅킹 시뮬레이션: 전체 파이프라인 3단계 시뮬",
    "전체 상태 초기화: 세션 + 이력 전체 삭제",
    "현재 세션 상태 실시간 표시",
]
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(6), Inches(3.5), debug_items, size=16)
add_phone_img(slide, os.path.join(BASE, "a_after.png"), Inches(8.5), Inches(1.3))
add_text_box(slide, Inches(8.5), Inches(5.1), PHONE_W, Inches(0.5),
             "디버그 패널 (앱 설정)", size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━ SLIDE 9-A: 회귀 테스트 12항목 본셋 ★필수1 전면 재작성 ━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "회귀 테스트 진행 상태 — 12항목 본셋")

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(12), Inches(0.5),
             "코드 경로 확인 대부분 완료. 실기 검증 및 단위 테스트 잔여.",
             size=18, bold=True)

# 헤더
hy = Inches(1.9)
add_box(slide, Inches(0.4), hy, Inches(0.6), Inches(0.42), NAVY, "#", 12)
add_box(slide, Inches(1.05), hy, Inches(4.5), Inches(0.42), NAVY, "검증 항목", 12)
add_box(slide, Inches(5.6), hy, Inches(3.8), Inches(0.42), NAVY, "검증 방법", 12)
add_box(slide, Inches(9.45), hy, Inches(3.5), Inches(0.42), NAVY, "현재 상태", 12)

results_9a = [
    ("#1", "same-call snooze 후 suppress 유지", "디버그 훅 + 실기 또는 단위 테스트", "코드 경로 확인 완료", BLUE),
    ("#2", "snooze 후 텔레뱅킹 윈도우 내 bank outgoing", "단위 테스트", "사양 정합 확정", BLUE),
    ("#3", "snooze 후 일반 outgoing 미방출", "디버그 훅 또는 단위 테스트", "대기", GRAY),
    ("#4", "텔레뱅킹 윈도우 만료 후 미방출", "단위 테스트 (clock 주입)", "대기", GRAY),
    ("#5", "통화 종료(IDLE) 후 snooze clear", "실기 (PR-α1) 또는 단위 테스트", "코드 경로 확인 완료", BLUE),
    ("#6", "CRITICAL 팝업 + 쿨다운 팝업 최상위", "디버그 훅", "대기", GRAY),
    ("#7", "팝업 dismiss 후 쿨다운 노출", "실기 풀 시퀀스", "코드 경로: 자동 연쇄 없음", ORANGE),
    ("#8", "팝업 뷰 상태 재배치", "—", "제외 (가시 상태 부재)", GRAY),
    ("#9", "ensureCriticalOnTop 중복 호출", "디버그 훅 또는 단위 테스트", "코드 확인 완료", BLUE),
    ("#10", "no-parent 무크래시", "단위 테스트", "코드 확인 완료", BLUE),
    ("#11", "dismiss race 무크래시", "단위 테스트", "코드 확인 완료", BLUE),
    ("#12", "off-main serialize", "단위 테스트", "코드 확인 완료", BLUE),
]

for i, (num, item, method, status, color) in enumerate(results_9a):
    y = Inches(2.42 + i * 0.4)
    add_text_box(slide, Inches(0.5), y, Inches(0.5), Inches(0.35), num, size=11, bold=True)
    add_text_box(slide, Inches(1.1), y, Inches(4.4), Inches(0.35), item, size=11)
    add_text_box(slide, Inches(5.65), y, Inches(3.7), Inches(0.35), method, size=11, color=GRAY)
    add_box(slide, Inches(9.45), y, Inches(3.5), Inches(0.35), color, status, 10)

# ━━━━━━━━━━━ SLIDE 9-B: 파생 점검 및 다음 단계 ★필수1 ━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "파생 점검 및 다음 단계")

# 파생 점검
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(12), Inches(0.5),
             "파생 점검 (12항목 본셋 외)", size=20, bold=True)
add_box(slide, Inches(0.8), Inches(2.2), Inches(1.0), Inches(0.45), BLUE, "#5b", 14)
add_text_box(slide, Inches(2.0), Inches(2.2), Inches(10), Inches(0.45),
             "snooze TTL 15분 만료 후 clear — 코드 경로 확인 완료, 단위 테스트 권장", size=16)

# 다음 단계
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(12), Inches(0.5),
             "다음 단계", size=20, bold=True)
next_steps = [
    "PR-α1: #5 단독 실기 검증 (IDLE → snooze clear)",
    "디버그 훅 PR: #1, #3, #6, #9 검증용 디버그 트리거 추가",
    "단위 테스트 PR: #2, #4, #10, #11, #12 자동 검증",
]
add_bullet_list(slide, Inches(0.8), Inches(3.9), Inches(8), Inches(2.0), next_steps, size=18)

# 별도 트랙
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(0.5),
             "별도 트랙", size=20, bold=True)
add_box(slide, Inches(0.8), Inches(6.1), Inches(12), Inches(0.7), PURPLE_BG,
        "C-3 시맨틱 정리: reset()의 session/snooze 동시 초기화 분리 — 12항목 회귀와 독립 트랙",
        font_size=15, font_color=NAVY)

# ━━━━━━━━━━━ SLIDE 10: 알려진 이슈 ★필수2,3 수정 ━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "알려진 이슈 및 제한 사항")

issues = [
    ("UsageStats 간헐적 조회 실패: AT-M140에서 권한 허용 상태에서도 raw=0 반환. "
     "권한 toggle off/on + 앱 재시작으로 복구."),
    # ★ 필수3: 조건 명시 (B안)
    ("\"안전 확인했어요\" 동작 범위: 활성 위험 세션이 있을 때만 복귀 동작. "
     "세션 없는 상태에서는 버튼이 반응하지 않음 (UX 개선 예정)."),
    # ★ 필수2: C-3 문구 수정
    ("C-3 상태 분리 미구현: reset()이 session과 snooze를 함께 초기화하는 시맨틱 혼재. "
     "단 RiskOverlayManager의 재등록 경로로 same-call snooze 기능은 보장됨. "
     "시맨틱 정리는 별도 티켓."),
    "ANSWER_PHONE_CALLS 권한: Manifest에 선언만 존재, 실제 미사용. 제거 예정.",
]

for i, issue in enumerate(issues):
    y = Inches(1.5 + i * 1.0)
    add_box(slide, Inches(0.6), y, Inches(0.4), Inches(0.4), ORANGE, str(i + 1), 14)
    add_text_box(slide, Inches(1.2), y, Inches(11.5), Inches(0.8), issue, size=14)

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(12), Inches(0.5),
             "문의처", size=22, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.4), Inches(12), Inches(0.5),
             "Mobile QA Team  |  jinwu@altech.kr", size=18, color=BLUE)

# ━━━━━━━━━━━ 저장 ━━━━━━━━━━━
prs.save(OUT)
print(f"PPT 생성 완료: {OUT}")
